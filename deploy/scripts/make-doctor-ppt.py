#!/usr/bin/env python3
"""CAMA Plus 신규 기능 소개 PPT — 의사 대상, 부드러운 톤."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

# --- Palette (soft medical, warm — not stiff corporate) ---
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
TEAL_DARK = RGBColor(0x1D, 0x7A, 0x6F)
TEAL_SOFT = RGBColor(0xD8, 0xF0, 0xEB)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
CORAL_SOFT = RGBColor(0xFA, 0xE5, 0xDE)
CREAM = RGBColor(0xF8, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x2C, 0x35, 0x3A)
MUTED = RGBColor(0x5F, 0x6B, 0x72)
LINE = RGBColor(0xE4, 0xDE, 0xD4)
GOLD = RGBColor(0xC9, 0xA2, 0x27)

OUT = Path(__file__).resolve().parents[2] / "docs" / "CAMA_Plus_신규기능_의사설명.pptx"
W, H = Inches(13.333), Inches(7.5)


def set_run(run, size=18, bold=False, color=INK, font="맑은 고딕"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    # East Asian
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        from lxml import etree

        ea = etree.SubElement(
            rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea"
        )
    ea.set("typeface", font)


def add_text(shape, lines, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, spacing=1.15):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=bold, color=color)


def add_para(tf, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, space_after=8):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    if tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1:
        p = tf.paragraphs[0]
    else:
        if tf.paragraphs[0].text == "" and len([x for x in tf.paragraphs]) == 1:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, left, top, width, height, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(sh, color)
    # softer corners
    try:
        sh.adjustments[0] = 0.1
    except Exception:
        pass
    return sh


def oval(slide, left, top, width, height, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    fill_shape(sh, color)
    return sh


def textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def bg(slide, color=CREAM):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill_shape(sh, color)
    # send to back is not easy; leave as first shape
    return sh


def footer(slide, page: str, note: str = "CAMA Plus · 의사 선생님을 위한 안내"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), W, Inches(0.35))
    fill_shape(bar, WHITE)
    bar.line.color.rgb = LINE
    tb = textbox(slide, Inches(0.5), Inches(7.18), Inches(10), Inches(0.28))
    add_text(tb, [note], size=10, color=MUTED)
    tb2 = textbox(slide, Inches(11.5), Inches(7.18), Inches(1.5), Inches(0.28))
    add_text(tb2, [page], size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def section_chip(slide, text, left=Inches(0.55), top=Inches(0.35)):
    chip = rect(slide, left, top, Inches(2.2), Inches(0.38), TEAL_SOFT)
    tb = textbox(slide, left, top + Inches(0.05), Inches(2.2), Inches(0.3))
    add_text(tb, [text], size=12, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)


def title_block(slide, title, subtitle=None):
    section_chip(slide, "CAMA Plus")
    tb = textbox(slide, Inches(0.55), Inches(0.85), Inches(12), Inches(0.7))
    add_text(tb, [title], size=32, bold=True, color=INK)
    if subtitle:
        tb2 = textbox(slide, Inches(0.55), Inches(1.5), Inches(12), Inches(0.45))
        add_text(tb2, [subtitle], size=16, color=MUTED)


def card(slide, left, top, width, height, fill=WHITE):
    return rect(slide, left, top, width, height, fill)


def card_title_body(slide, left, top, width, height, title, body_lines, accent=TEAL):
    card(slide, left, top, width, height, WHITE)
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height
    )
    fill_shape(accent_bar, accent)
    tb = textbox(slide, left + Inches(0.3), top + Inches(0.22), width - Inches(0.45), Inches(0.4))
    add_text(tb, [title], size=16, bold=True, color=INK)
    tb2 = textbox(
        slide,
        left + Inches(0.3),
        top + Inches(0.65),
        width - Inches(0.45),
        height - Inches(0.85),
    )
    tf = tb2.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        set_run(run, size=13, color=MUTED)


def make():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    pages = []

    def new():
        s = prs.slides.add_slide(blank)
        bg(s, CREAM)
        pages.append(s)
        return s

    # ========== 1 Cover ==========
    s = new()
    # decorative blobs
    oval(s, Inches(-1.2), Inches(-1.5), Inches(5), Inches(5), TEAL_SOFT)
    oval(s, Inches(10.5), Inches(4.2), Inches(4), Inches(4), CORAL_SOFT)
    rect(s, Inches(1.2), Inches(1.6), Inches(10.9), Inches(4.4), WHITE)
    chip = rect(s, Inches(1.6), Inches(2.0), Inches(3.0), Inches(0.4), TEAL_SOFT)
    tb = textbox(s, Inches(1.6), Inches(2.05), Inches(3.0), Inches(0.35))
    add_text(tb, ["의사 선생님께 드리는 안내"], size=12, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

    tb = textbox(s, Inches(1.6), Inches(2.6), Inches(10), Inches(1.2))
    add_text(tb, ["진료가 조금 더 편해지도록,"], size=20, color=MUTED)
    tb = textbox(s, Inches(1.6), Inches(3.15), Inches(10), Inches(1.0))
    add_text(tb, ["CAMA Plus에 새로 생긴 기능들"], size=36, bold=True, color=INK)

    tb = textbox(s, Inches(1.6), Inches(4.4), Inches(10), Inches(0.8))
    add_text(
        tb,
        ["환자가 미리 준비하고 · 의사가 진료 전에 바로 볼 수 있게"],
        size=16,
        color=MUTED,
    )
    tb = textbox(s, Inches(1.6), Inches(5.2), Inches(10), Inches(0.4))
    add_text(tb, ["로그인 · 걸음·심박 · 문의(말로 작성) · 태블릿 연동 · 환자 관리"], size=13, color=TEAL_DARK)

    # ========== 2 Agenda ==========
    s = new()
    title_block(s, "오늘 함께 보실 다섯 가지", "기술 용어보다는, 진료 현장에서 어떻게 쓰이는지에 맞춰 정리했습니다.")
    items = [
        ("01", "로그인·계정", "환자가 스스로 가입하고, 비밀번호를 찾고 바꿀 수 있어요"),
        ("02", "생체 정보", "걸음수 등이 앱 사용 중에 자연스럽게 쌓여요"),
        ("03", "의사 문의", "키보드로도, 말로도 진료 전 질문을 남길 수 있어요"),
        ("04", "태블릿 연동", "환자 폰 → 의사 태블릿으로 바로 전달됩니다"),
        ("05", "환자 관리", "병원·관리 화면에서 환자·서비스를 살펴볼 수 있어요"),
    ]
    for i, (num, title, desc) in enumerate(items):
        top = Inches(2.15) + Inches(i * 0.9)
        card(s, Inches(0.55), top, Inches(12.2), Inches(0.78), WHITE)
        num_box = rect(s, Inches(0.75), top + Inches(0.18), Inches(0.7), Inches(0.42), TEAL_SOFT)
        tb = textbox(s, Inches(0.75), top + Inches(0.22), Inches(0.7), Inches(0.35))
        add_text(tb, [num], size=14, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
        tb = textbox(s, Inches(1.7), top + Inches(0.12), Inches(4), Inches(0.35))
        add_text(tb, [title], size=18, bold=True, color=INK)
        tb = textbox(s, Inches(1.7), top + Inches(0.42), Inches(10.5), Inches(0.3))
        add_text(tb, [desc], size=13, color=MUTED)

    # ========== 3 Story ==========
    s = new()
    title_block(
        s,
        "한 명의 환자, 한 번의 진료로 보는 이야기",
        "기능 나열이 아니라, 진료 전까지의 흐름으로 따라가 볼게요.",
    )
    steps = [
        ("가입·로그인", "계정을 만들고\n들어옵니다"),
        ("걸음이 쌓임", "앱을 쓰는 동안\n활동량이 저장됩니다"),
        ("문의 작성", "말로도 질문을\n미리 남깁니다"),
        ("자료 전송", "QR로 태블릿에\n가볍게 보냅니다"),
        ("의사 확인", "건강·문의를\n한눈에 봅니다"),
    ]
    for i, (t, d) in enumerate(steps):
        left = Inches(0.45) + Inches(i * 2.55)
        card(s, left, Inches(2.5), Inches(2.35), Inches(3.2), WHITE)
        circ = oval(s, left + Inches(0.75), Inches(2.75), Inches(0.85), Inches(0.85), TEAL_SOFT)
        tb = textbox(s, left + Inches(0.75), Inches(2.95), Inches(0.85), Inches(0.5))
        add_text(tb, [str(i + 1)], size=20, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
        tb = textbox(s, left + Inches(0.15), Inches(3.8), Inches(2.05), Inches(0.5))
        add_text(tb, [t], size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
        tb = textbox(s, left + Inches(0.15), Inches(4.4), Inches(2.05), Inches(1.0))
        add_text(tb, d.split("\n"), size=12, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 4:
            tb = textbox(s, left + Inches(2.2), Inches(3.7), Inches(0.4), Inches(0.4))
            add_text(tb, ["→"], size=20, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

    # ========== 4 Before After ==========
    s = new()
    title_block(s, "무엇이 달라졌을까요?", "이전과 지금을 가볍게 비교해 보았습니다.")
    rows = [
        ("계정", "찾기·수정이 번거로웠어요", "찾기·초기화, 이메일, 내정보·비번 변경"),
        ("생체정보", "특정 화면에서만 걸음이 저장됐어요", "앱을 켜면 자연스럽게 동기화"),
        ("의사 문의", "미리 질문을 남기기 어려웠어요", "텍스트 + 말로 작성"),
        ("진료 현장", "자료가 흩어져 있었어요", "환자 폰 → 의사 태블릿으로 바로"),
        ("환자 관리", "운영 화면이 들쭉날쭉했어요", "모니터링·서비스·APK 등 정비"),
    ]
    # headers
    card(s, Inches(0.55), Inches(2.15), Inches(12.2), Inches(0.55), TEAL)
    for left, text in [
        (Inches(0.7), "영역"),
        (Inches(2.8), "이전"),
        (Inches(7.5), "지금은"),
    ]:
        tb = textbox(s, left, Inches(2.25), Inches(4), Inches(0.4))
        add_text(tb, [text], size=14, bold=True, color=WHITE)

    for i, (a, b, c) in enumerate(rows):
        top = Inches(2.8) + Inches(i * 0.75)
        fill = WHITE if i % 2 == 0 else TEAL_SOFT
        card(s, Inches(0.55), top, Inches(12.2), Inches(0.68), fill)
        tb = textbox(s, Inches(0.7), top + Inches(0.18), Inches(2), Inches(0.4))
        add_text(tb, [a], size=14, bold=True, color=TEAL_DARK)
        tb = textbox(s, Inches(2.8), top + Inches(0.18), Inches(4.4), Inches(0.4))
        add_text(tb, [b], size=13, color=MUTED)
        tb = textbox(s, Inches(7.5), top + Inches(0.18), Inches(5), Inches(0.4))
        add_text(tb, [c], size=13, bold=True, color=INK)

    # ========== 5 Section Account ==========
    s = new()
    oval(s, Inches(9), Inches(-1), Inches(5), Inches(5), TEAL_SOFT)
    section_chip(s, "첫 번째")
    tb = textbox(s, Inches(0.7), Inches(2.5), Inches(11), Inches(1))
    add_text(tb, ["01  로그인 · 계정"], size=40, bold=True, color=INK)
    tb = textbox(s, Inches(0.7), Inches(3.6), Inches(11), Inches(1))
    add_text(
        tb,
        ["환자가 계정을 스스로 관리할 수 있게 됐어요.", "선생님께서는 ‘접속이 잘 되나?’만 가볍게 알아두셔도 충분합니다."],
        size=18,
        color=MUTED,
    )

    # ========== 6 Account features ==========
    s = new()
    title_block(s, "계정이 이렇게 편해졌어요", "환자가 앱에서 직접 할 수 있는 일들입니다.")
    feats = [
        ("회원가입 · 로그인", "아이디와 비밀번호로\n서비스에 들어옵니다"),
        ("아이디 찾기", "이름과 전화번호로\n아이디를 확인합니다"),
        ("비밀번호 초기화", "임시 비밀번호를 받고\n다시 로그인합니다"),
        ("이메일 안내", "이메일이 있으면\n임시 비번·변경 안내 메일"),
        ("내 상세정보", "이름·전화·이메일·\n성별·생년월일 확인"),
        ("프로필 · 비번 변경", "정보 수정과 비밀번호\n변경을 같은 화면에서"),
    ]
    for i, (t, d) in enumerate(feats):
        col, row = i % 3, i // 3
        left = Inches(0.55) + Inches(col * 4.2)
        top = Inches(2.2) + Inches(row * 2.2)
        card_title_body(
            s,
            left,
            top,
            Inches(3.95),
            Inches(1.95),
            t,
            d.split("\n"),
            accent=TEAL if row == 0 else CORAL,
        )

    # ========== 7 Account tip ==========
    s = new()
    title_block(s, "의사 선생님이 알아두면 좋은 점", "환자에게 안내하실 때 참고해 주세요.")
    tips = [
        ("비밀번호를 잊었을 때", "앱의 ‘계정 찾기’에서 초기화할 수 있어요. 이메일을 등록해 두면 메일로도 안내됩니다."),
        ("비밀번호를 바꿀 때", "내 상세정보 → 수정에서 현재 비밀번호를 확인한 뒤 새 비밀번호로 바꿉니다. (일반 웹사이트와 같아요)"),
        ("이메일 등록을 권장", "초기화·변경 안내를 메일로 받을 수 있어, 진료 중 ‘접속이 안 돼요’ 문의가 줄어들 수 있어요."),
    ]
    for i, (t, d) in enumerate(tips):
        top = Inches(2.2) + Inches(i * 1.4)
        card(s, Inches(0.55), top, Inches(12.2), Inches(1.2), WHITE)
        badge = rect(s, Inches(0.8), top + Inches(0.35), Inches(0.55), Inches(0.55), CORAL_SOFT)
        tb = textbox(s, Inches(0.8), top + Inches(0.45), Inches(0.55), Inches(0.4))
        add_text(tb, [str(i + 1)], size=16, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
        tb = textbox(s, Inches(1.6), top + Inches(0.25), Inches(10.5), Inches(0.35))
        add_text(tb, [t], size=18, bold=True, color=INK)
        tb = textbox(s, Inches(1.6), top + Inches(0.65), Inches(10.5), Inches(0.4))
        add_text(tb, [d], size=14, color=MUTED)

    # ========== 8 Section Vital ==========
    s = new()
    oval(s, Inches(-1), Inches(4), Inches(4.5), Inches(4.5), CORAL_SOFT)
    section_chip(s, "두 번째")
    tb = textbox(s, Inches(0.7), Inches(2.5), Inches(11), Inches(1))
    add_text(tb, ["02  생체 정보 수집"], size=40, bold=True, color=INK)
    tb = textbox(s, Inches(0.7), Inches(3.6), Inches(11), Inches(1.2))
    add_text(
        tb,
        ["걸음수가 앱을 쓰는 동안 알아서 쌓입니다.", "심박 등 바이탈도 같은 ‘건강 데이터’ 흐름으로 확장되고 있어요."],
        size=18,
        color=MUTED,
    )

    # ========== 9 Steps ==========
    s = new()
    title_block(s, "걸음수가 알아서 쌓여요", "환자가 따로 ‘전송’ 버튼을 누르지 않아도 됩니다.")
    moments = [
        ("로그인 후", "앱에 들어오면\n오늘 걸음을 맞춰 둡니다"),
        ("앱을 다시 켤 때", "포그라운드로 돌아오면\n다시 한 번 맞춥니다"),
        ("자료 전송 직전", "의사 태블릿으로 보내기 전\n최신 값으로 갱신합니다"),
        ("마이페이지", "걸음수 메뉴에서\n이력을 확인할 수 있어요"),
    ]
    for i, (t, d) in enumerate(moments):
        left = Inches(0.5) + Inches(i * 3.2)
        card(s, left, Inches(2.4), Inches(3.0), Inches(3.5), WHITE)
        top_band = rect(s, left, Inches(2.4), Inches(3.0), Inches(0.7), TEAL if i % 2 == 0 else CORAL)
        tb = textbox(s, left, Inches(2.55), Inches(3.0), Inches(0.45))
        add_text(tb, [t], size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb = textbox(s, left + Inches(0.2), Inches(3.5), Inches(2.6), Inches(1.8))
        add_text(tb, d.split("\n"), size=15, color=MUTED, align=PP_ALIGN.CENTER)

    # ========== 10 Vital note ==========
    s = new()
    title_block(s, "심박 등 바이탈은 이렇게 이해해 주세요", "걸음과 같은 ‘건강 데이터’ 축입니다.")
    card(s, Inches(0.55), Inches(2.3), Inches(6.0), Inches(4.0), WHITE)
    tb = textbox(s, Inches(0.9), Inches(2.6), Inches(5.3), Inches(0.5))
    add_text(tb, ["포함되는 개념"], size=18, bold=True, color=INK)
    for i, line in enumerate(
        ["심박수", "혈압", "산소포화도", "체온 · 호흡 등", "서버에 이력으로 쌓이고, 태블릿 차트와도 연결"]
    ):
        tb = textbox(s, Inches(0.9), Inches(3.3) + Inches(i * 0.5), Inches(5.3), Inches(0.45))
        add_text(tb, [f"•  {line}"], size=15, color=MUTED)

    card(s, Inches(6.85), Inches(2.3), Inches(5.9), Inches(4.0), TEAL_SOFT)
    tb = textbox(s, Inches(7.2), Inches(2.8), Inches(5.2), Inches(2.5))
    add_text(
        tb,
        [
            "진료실에서의 의미",
            "",
            "환자가 평소에 움직인 양과",
            "건강 기록이 태블릿에 모이면,",
            "짧은 진료 시간에도",
            "‘최근 상태’를 빠르게 파악하실 수 있어요.",
        ],
        size=16,
        color=TEAL_DARK,
        align=PP_ALIGN.LEFT,
    )

    # ========== 11 Section Inquiry ==========
    s = new()
    oval(s, Inches(10), Inches(-0.5), Inches(4), Inches(4), TEAL_SOFT)
    section_chip(s, "세 번째")
    tb = textbox(s, Inches(0.7), Inches(2.5), Inches(11), Inches(1))
    add_text(tb, ["03  의사 문의 작성"], size=40, bold=True, color=INK)
    tb = textbox(s, Inches(0.7), Inches(3.6), Inches(11), Inches(1.2))
    add_text(
        tb,
        ["진찰 전에 환자가 질문을 미리 남겨 둡니다.", "타이핑도 되고, 말로 말해도 글자로 들어갑니다."],
        size=18,
        color=MUTED,
    )

    # ========== 12 Inquiry detail ==========
    s = new()
    title_block(s, "「진찰시 문의사항」은 이런 기능이에요", "마이페이지에서 환자가 직접 작성합니다.")
    card_title_body(
        s,
        Inches(0.55),
        Inches(2.2),
        Inches(6.0),
        Inches(4.2),
        "무엇을 할 수 있나요?",
        [
            "• 문의 작성 · 수정 · 삭제",
            "• 제목과 본문을 남겨 둡니다",
            "• 진료 전에 ‘물어볼 말’을 정리",
            "• 나중에 태블릿으로 전달",
            "",
            "환자가 말로 준비한 질문이",
            "진료실에서 바로 보입니다.",
        ],
        accent=TEAL,
    )
    card(s, Inches(6.85), Inches(2.2), Inches(5.9), Inches(4.2), CORAL_SOFT)
    tb = textbox(s, Inches(7.2), Inches(2.5), Inches(5.2), Inches(0.5))
    add_text(tb, ["입력 방식"], size=18, bold=True, color=INK)
    # two mini cards
    card(s, Inches(7.2), Inches(3.2), Inches(5.2), Inches(1.3), WHITE)
    tb = textbox(s, Inches(7.45), Inches(3.4), Inches(4.8), Inches(1.0))
    add_text(tb, ["키보드로 입력", "평소처럼 글자를 적습니다"], size=15, color=MUTED)
    card(s, Inches(7.2), Inches(4.7), Inches(5.2), Inches(1.4), WHITE)
    tb = textbox(s, Inches(7.45), Inches(4.9), Inches(4.8), Inches(1.1))
    add_text(
        tb,
        ["말로 입력 (마이크)", "말하면 → 한국어로 글자가 붙습니다"],
        size=15,
        color=MUTED,
    )

    # ========== 13 Voice highlight ==========
    s = new()
    title_block(s, "말로 입력해도, 결국 ‘텍스트’로 남아요", "음성 파일을 따로 저장하는 방식이 아닙니다.")
    card(s, Inches(1.5), Inches(2.4), Inches(10.3), Inches(3.8), WHITE)
    flow = [
        ("말하기", "환자가 마이크에\n증상을 말해요"),
        ("글자로", "앱이 한국어로\n인식해 붙여 줘요"),
        ("저장", "문의 본문으로\n서버에 남아요"),
        ("의사 확인", "태블릿에서\n읽기 쉬운 글자로"),
    ]
    for i, (t, d) in enumerate(flow):
        left = Inches(1.8) + Inches(i * 2.5)
        circ = oval(s, left + Inches(0.55), Inches(2.8), Inches(0.9), Inches(0.9), TEAL_SOFT)
        tb = textbox(s, left + Inches(0.55), Inches(3.0), Inches(0.9), Inches(0.5))
        add_text(tb, [str(i + 1)], size=18, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
        tb = textbox(s, left, Inches(3.9), Inches(2.1), Inches(0.4))
        add_text(tb, [t], size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
        tb = textbox(s, left, Inches(4.4), Inches(2.1), Inches(1.0))
        add_text(tb, d.split("\n"), size=13, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 3:
            tb = textbox(s, left + Inches(2.05), Inches(3.0), Inches(0.4), Inches(0.5))
            add_text(tb, ["→"], size=18, bold=True, color=CORAL)

    # ========== 14 Section Tablet ==========
    s = new()
    oval(s, Inches(-1.5), Inches(-1), Inches(5), Inches(5), CORAL_SOFT)
    section_chip(s, "네 번째")
    tb = textbox(s, Inches(0.7), Inches(2.5), Inches(11), Inches(1))
    add_text(tb, ["04  의사 태블릿 연동"], size=40, bold=True, color=INK)
    tb = textbox(s, Inches(0.7), Inches(3.6), Inches(11), Inches(1.2))
    add_text(
        tb,
        ["환자 스마트폰에서 의사 태블릿으로, 바로 자료를 넘깁니다.", "QR을 스캔하고 보내면, 건강 데이터와 문의가 화면에 나타납니다."],
        size=18,
        color=MUTED,
    )

    # ========== 15 Tablet flow ==========
    s = new()
    title_block(s, "진료실에서의 짧은 흐름", "복잡한 설정 없이, 이렇게만 하시면 됩니다.")
    steps = [
        ("1", "태블릿", "QR 대기 화면을\n띄워 둡니다"),
        ("2", "환자 앱", "「의사앱 자료전송」에서\nQR을 스캔합니다"),
        ("3", "미리보기", "보낼 걸음·문의를\n확인합니다"),
        ("4", "전송", "태블릿으로\n자료를 보냅니다"),
        ("5", "확인", "건강·문의 탭에서\n내용을 봅니다"),
    ]
    for i, (n, t, d) in enumerate(steps):
        left = Inches(0.4) + Inches(i * 2.55)
        card(s, left, Inches(2.4), Inches(2.4), Inches(3.8), WHITE)
        band = rect(s, left, Inches(2.4), Inches(2.4), Inches(0.9), TEAL if i != 4 else CORAL)
        tb = textbox(s, left, Inches(2.5), Inches(2.4), Inches(0.35))
        add_text(tb, [n], size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb = textbox(s, left, Inches(2.85), Inches(2.4), Inches(0.4))
        add_text(tb, [t], size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb = textbox(s, left + Inches(0.15), Inches(3.6), Inches(2.1), Inches(2.0))
        add_text(tb, d.split("\n"), size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # ========== 16 Tablet tabs ==========
    s = new()
    title_block(s, "태블릿에서 이렇게 보입니다", "두 개의 탭으로 나눠 두었어요.")
    card(s, Inches(0.55), Inches(2.3), Inches(6.0), Inches(4.1), WHITE)
    top = rect(s, Inches(0.55), Inches(2.3), Inches(6.0), Inches(0.8), TEAL)
    tb = textbox(s, Inches(0.55), Inches(2.5), Inches(6.0), Inches(0.45))
    add_text(tb, ["건강 데이터"], size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, line in enumerate(
        [
            "환자 이름 · 기본 정보",
            "걸음수와 최근 이력",
            "심박 등 바이탈 차트 영역",
            "짧은 시간에 ‘최근 활동’을 파악",
        ]
    ):
        tb = textbox(s, Inches(1.0), Inches(3.5) + Inches(i * 0.55), Inches(5), Inches(0.45))
        add_text(tb, [f"•  {line}"], size=16, color=MUTED)

    card(s, Inches(6.85), Inches(2.3), Inches(5.9), Inches(4.1), WHITE)
    top = rect(s, Inches(6.85), Inches(2.3), Inches(5.9), Inches(0.8), CORAL)
    tb = textbox(s, Inches(6.85), Inches(2.5), Inches(5.9), Inches(0.45))
    add_text(tb, ["문의사항"], size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, line in enumerate(
        [
            "환자가 미리 남긴 질문 목록",
            "제목 · 본문 상세 보기",
            "말로 작성한 내용도 텍스트로",
            "문진 전에 ‘오늘 궁금한 점’ 확인",
        ]
    ):
        tb = textbox(s, Inches(7.3), Inches(3.5) + Inches(i * 0.55), Inches(5), Inches(0.45))
        add_text(tb, [f"•  {line}"], size=16, color=MUTED)

    # ========== 17 Section Admin ==========
    s = new()
    oval(s, Inches(9.5), Inches(3.5), Inches(5), Inches(5), TEAL_SOFT)
    section_chip(s, "다섯 번째")
    tb = textbox(s, Inches(0.7), Inches(2.5), Inches(11), Inches(1))
    add_text(tb, ["05  환자 관리앱"], size=40, bold=True, color=INK)
    tb = textbox(s, Inches(0.7), Inches(3.6), Inches(11), Inches(1.2))
    add_text(
        tb,
        ["진료 태블릿과 별도로, 병원·관리 웹에서 환자를 살펴보는 기능입니다.", "운영·모니터링이 필요할 때 쓰시는 화면이에요."],
        size=18,
        color=MUTED,
    )

    # ========== 18 Admin features ==========
    s = new()
    title_block(s, "관리자 · 의사 웹에 담긴 기능", "역할에 따라 보시면 됩니다.")
    card(s, Inches(0.55), Inches(2.15), Inches(6.0), Inches(4.3), WHITE)
    band = rect(s, Inches(0.55), Inches(2.15), Inches(6.0), Inches(0.65), TEAL)
    tb = textbox(s, Inches(0.55), Inches(2.3), Inches(6.0), Inches(0.4))
    add_text(tb, ["슈퍼 관리자"], size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, (t, d) in enumerate(
        [
            ("환자 모니터링", "환자 목록 · 유형 · 상세 연계"),
            ("코칭 모니터링", "환자별 코칭 진행 현황"),
            ("월평가 지표", "월별 지표 조회 · 엑셀"),
            ("APK 관리", "환자 앱 버전 업로드 · 배포"),
            ("치료정보 사용현황", "치료정보 사용 현황 확인"),
        ]
    ):
        tb = textbox(s, Inches(0.9), Inches(3.05) + Inches(i * 0.6), Inches(5.3), Inches(0.55))
        add_text(tb, [f"{t}  —  {d}"], size=13, color=MUTED)

    card(s, Inches(6.85), Inches(2.15), Inches(5.9), Inches(4.3), WHITE)
    band = rect(s, Inches(6.85), Inches(2.15), Inches(5.9), Inches(0.65), CORAL)
    tb = textbox(s, Inches(6.85), Inches(2.3), Inches(5.9), Inches(0.4))
    add_text(tb, ["의사 웹"], size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, (t, d) in enumerate(
        [
            ("환자 관리", "병원 환자 목록 · 모니터링"),
            ("서비스 신청", "신청 목록 · 승인 · 반려"),
            ("연동 안정화", "화면과 서버 호출이 더 안정적으로"),
            ("", ""),
            ("포인트", "현장 문의는 태블릿, 운영은 웹"),
        ]
    ):
        if not t:
            continue
        tb = textbox(s, Inches(7.2), Inches(3.05) + Inches(i * 0.65), Inches(5.2), Inches(0.55))
        add_text(tb, [f"{t}  —  {d}"], size=13, color=MUTED)

    # ========== 19 Roles ==========
    s = new()
    title_block(s, "누가 무엇을 볼까요?", "헷갈리지 않게, 역할만 정리했습니다.")
    roles = [
        ("환자", "계정 · 걸음 · 문의 작성 · 자료 전송", TEAL),
        ("의사 (태블릿)", "진료 직전 건강 데이터 · 문의사항", CORAL),
        ("의사 (웹)", "환자 목록 · 서비스 신청 처리", TEAL_DARK),
        ("관리자", "모니터링 · APK · 지표 · 운영", GOLD),
    ]
    for i, (t, d, c) in enumerate(roles):
        left = Inches(0.5) + Inches(i * 3.2)
        card(s, left, Inches(2.5), Inches(3.0), Inches(3.5), WHITE)
        top = rect(s, left, Inches(2.5), Inches(3.0), Inches(1.0), c)
        tb = textbox(s, left, Inches(2.75), Inches(3.0), Inches(0.55))
        add_text(tb, [t], size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb = textbox(s, left + Inches(0.2), Inches(3.9), Inches(2.6), Inches(1.6))
        add_text(tb, [d], size=15, color=MUTED, align=PP_ALIGN.CENTER)

    # ========== 20 Recap ==========
    s = new()
    title_block(s, "한 장으로 다시 보기", "오늘 말씀드린 핵심만 모았습니다.")
    recap = [
        ("01", "계정", "환자가 스스로 찾고, 바꾸고, 이메일도 등록"),
        ("02", "생체정보", "걸음이 자동으로 쌓이고, 바이탈로 확장"),
        ("03", "문의", "키보드 + 말로 진료 전 질문 준비"),
        ("04", "태블릿", "QR로 보내면 건강·문의가 바로 표시"),
        ("05", "관리", "웹에서 환자·서비스·앱을 운영"),
    ]
    for i, (n, t, d) in enumerate(recap):
        top = Inches(2.15) + Inches(i * 0.9)
        card(s, Inches(0.55), top, Inches(12.2), Inches(0.78), WHITE)
        badge = rect(s, Inches(0.8), top + Inches(0.18), Inches(0.7), Inches(0.42), TEAL_SOFT)
        tb = textbox(s, Inches(0.8), top + Inches(0.22), Inches(0.7), Inches(0.35))
        add_text(tb, [n], size=13, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
        tb = textbox(s, Inches(1.75), top + Inches(0.22), Inches(2.2), Inches(0.4))
        add_text(tb, [t], size=16, bold=True, color=INK)
        tb = textbox(s, Inches(4.1), top + Inches(0.22), Inches(8.2), Inches(0.4))
        add_text(tb, [d], size=15, color=MUTED)

    # ========== 21 Closing ==========
    s = new()
    oval(s, Inches(-1), Inches(-1), Inches(4.5), Inches(4.5), TEAL_SOFT)
    oval(s, Inches(10), Inches(4), Inches(4.5), Inches(4.5), CORAL_SOFT)
    card(s, Inches(1.8), Inches(1.8), Inches(9.7), Inches(4.0), WHITE)
    tb = textbox(s, Inches(2.2), Inches(2.2), Inches(9), Inches(0.5))
    add_text(tb, ["감사합니다"], size=36, bold=True, color=INK, align=PP_ALIGN.CENTER)
    tb = textbox(s, Inches(2.2), Inches(3.1), Inches(9), Inches(1.2))
    add_text(
        tb,
        [
            "환자가 미리 준비하고,",
            "선생님께서는 진료 전에 가볍게 확인하시면 됩니다.",
        ],
        size=18,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    tb = textbox(s, Inches(2.2), Inches(4.5), Inches(9), Inches(0.6))
    add_text(
        tb,
        ["궁금하신 점 있으시면 편하게 말씀해 주세요."],
        size=15,
        color=TEAL_DARK,
        align=PP_ALIGN.CENTER,
    )

    # footers (skip cover & closing decorative)
    total = len(pages)
    for i, slide in enumerate(pages):
        if i == 0 or i == total - 1:
            continue
        footer(slide, f"{i}/{total - 1}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {total}")


if __name__ == "__main__":
    make()
